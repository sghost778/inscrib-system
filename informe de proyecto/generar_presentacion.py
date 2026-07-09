#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Genera la presentación del proyecto INSCRIB SYSTEM en formato .pptx nativo."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT_DIR = r'D:\Usuarios\MCAMPOS\Desktop\INSCRIB SYSTEM\INSCRIB SYSTEM\informe de proyecto'
OUTPUT = os.path.join(OUTPUT_DIR, 'PRESENTACION_PROYECTO_ESCOLA.pptx')

# Colores institucionales
AZUL = RGBColor(0x00, 0xB6, 0x89)      # verde/teal de la app
AZUL_OSC = RGBColor(0x0B, 0x3D, 0x2E)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS = RGBColor(0x33, 0x33, 0x33)
GRIS_CLARO = RGBColor(0xEE, 0xF4, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=GRIS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=GRIS):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run(); r.text = '•  ' + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = 'Calibri'
    return tb


def title_bar(slide, title):
    add_rect(slide, 0, 0, SW, Inches(1.2), AZUL)
    add_text(slide, Inches(0.5), 0, SW - Inches(1), Inches(1.2), title,
             size=30, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)


# ── Diapositiva 1: Portada ──────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, AZUL)
add_text(s, Inches(0.8), Inches(2.0), SW - Inches(1.6), Inches(1.0),
         'INSCRIB SYSTEM', size=54, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.1), SW - Inches(1.6), Inches(0.8),
         'Sistema de Gestión de Inscripciones y Administración Escolar',
         size=24, bold=True, color=GRIS_CLARO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.3), SW - Inches(1.6), Inches(0.6),
         'U.E. Dr. José Manuel Cova Maza - Puerto Ordaz, Estado Bolívar',
         size=18, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.1), SW - Inches(1.6), Inches(0.6),
         'Marcelo Fernando Campos Anacona (V-32.062.637)  |  UPTJAA  |  Pasantías Profesionales 2026',
         size=14, color=GRIS_CLARO, align=PP_ALIGN.CENTER)

# ── Diapositiva 2: El Problema ─────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'El Problema')
add_bullets(s, Inches(0.7), Inches(1.6), SW - Inches(1.4), Inches(5), [
    'Procesos de inscripción y matrícula manuales, basados en hojas de cálculo.',
    'Falta de trazabilidad de los cambios sobre los expedientes escolares.',
    'Demoras en la atención a representantes y en la generación de documentos.',
    'Información institucional no publicada ni actualizada en la web.',
    'Riesgo de pérdida de información por falta de respaldos y auditoría.',
], size=20)

# ── Diapositiva 3: Objetivos ───────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Objetivos')
add_text(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(0.5),
         'Objetivo General', size=22, bold=True, color=AZUL)
add_text(s, Inches(0.7), Inches(2.1), SW - Inches(1.4), Inches(1.0),
         'Desarrollar un sistema web que automatice las inscripciones, la gestión de '
         'estudiantes y representantes, y la administración del sitio web institucional.',
         size=18)
add_text(s, Inches(0.7), Inches(3.3), SW - Inches(1.4), Inches(0.5),
         'Objetivos Específicos', size=22, bold=True, color=AZUL)
add_bullets(s, Inches(0.7), Inches(3.9), SW - Inches(1.4), Inches(3), [
    'Diagnosticar los procesos de administración escolar.',
    'Diseñar la arquitectura con Flask y SQLAlchemy.',
    'Desarrollar módulos de estudiantes, representantes y matrícula.',
    'Desarrollar el sitio web público y la seguridad/auditoría.',
], size=18)

# ── Diapositiva 4: Arquitectura ────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Arquitectura del Sistema')


def add_table(slide, x, y, w, headers, rows, col_w=None):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    gtab = slide.shapes.add_table(rows_n, cols_n, x, y, w, Inches(0.5 * rows_n)).table
    for j, h in enumerate(headers):
        c = gtab.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = AZUL
        c.text = h
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = BLANCO; p.runs[0].font.name = 'Calibri'
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gtab.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = GRIS_CLARO if i % 2 else BLANCO
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = GRIS
            p.runs[0].font.name = 'Calibri'
    return gtab


add_table(s, Inches(0.7), Inches(1.6), SW - Inches(1.4),
          ['Capa', 'Tecnología'],
          [['Presentación', 'HTML5, CSS3, JavaScript, Font Awesome, PWA (manifest/sw.js)'],
           ['Lógica', 'Flask 3.1.1, Blueprints, SQLAlchemy, bcrypt, Talisman, Limiter'],
           ['Datos', 'SQLite (test.db) - 20 tablas']])

# ── Diapositiva 5: Módulos ─────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Módulos del Sistema')
add_bullets(s, Inches(0.7), Inches(1.6), SW - Inches(1.4), Inches(5.3), [
    'Autenticación / Login (bcrypt + Talisman + Limiter).',
    'Gestión de Usuarios (roles admin / secretario).',
    'Año Escolar y Grados.',
    'Estudiantes y Representantes (registro, listado, consulta).',
    'Matrícula / Inscripción.',
    'Plantillas de Documentos (python-docx).',
    'Administración del Sitio Web (Noticias, Programas, Galería, Config).',
    'Seguridad y Auditoría (REGISTRO_AUDITORIA).',
    'Portal del Representante (sitio público).',
], size=17)

# ── Diapositiva 6: Tecnologías ─────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Tecnologías Utilizadas')
add_table(s, Inches(0.7), Inches(1.6), SW - Inches(1.4),
          ['Tecnología', 'Versión', 'Propósito'],
          [['Flask', '3.1.1', 'Backend web (Python 3.13)'],
           ['SQLAlchemy', '2.0', 'ORM'],
           ['Flask-Talisman', '-', 'Cabeceras de seguridad / CSP'],
           ['Flask-Limiter', '-', 'Rate limiting'],
           ['bcrypt', '-', 'Hash de contraseñas'],
           ['python-docx', '-', 'Generación de documentos'],
           ['SQLite', '3.x', 'Base de datos'],
           ['PWA', '-', 'App instalable']])

# ── Diapositiva 7: Resultados ──────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Resultados')
add_bullets(s, Inches(0.7), Inches(1.6), SW - Inches(1.4), Inches(5.3), [
    'Sistema implementado con 9 módulos funcionales.',
    '20 tablas de base de datos modelando la gestión escolar.',
    'Usuarios por defecto (admin/secretario) creados en create_app.',
    'REGISTRO_AUDITORIA con clave INTEGER AUTOINCREMENT.',
    'INSCRIPCION con estado, fecha_retiro, lapso_registro y motivo_retiro.',
    'Sitio público en puerto 5002 y panel admin en puerto 5001.',
], size=18)

# ── Diapositiva 8: Conclusiones ────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLANCO)
title_bar(s, 'Conclusiones y Recomendaciones')
add_bullets(s, Inches(0.7), Inches(1.6), SW - Inches(1.4), Inches(5.3), [
    'El sistema cumple los objetivos y mejora la eficiencia administrativa.',
    'La seguridad (bcrypt, Talisman, Limiter) y la auditoría fortalecen la integridad.',
    'Recomendación: respaldos automáticos y migración a un motor más robusto en producción.',
    'Recomendación: ampliar el portal del representante y agregar reportes.',
], size=20)

prs.save(OUTPUT)
print(f'Presentación guardada: {OUTPUT}')
print(f'Tamaño: {os.path.getsize(OUTPUT) / 1024:.1f} KB')
